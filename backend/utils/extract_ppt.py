from pptx import Presentation
import sys

ppt_path = r"D:\Christ University\MSc\5MDS\ProjectFinance\PortifolioBuilder\Finmentor (1).pptx"

try:
    prs = Presentation(ppt_path)
    with open("ppt_output.txt", "w", encoding="utf-8") as f:
        f.write(f"Total slides: {len(prs.slides)}\n\n")
        f.write("="*80 + "\n")
        
        for i, slide in enumerate(prs.slides):
            f.write(f"\n--- SLIDE {i+1} ---\n")
            f.write("-" * 80 + "\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    f.write(shape.text + "\n")
            
            # Check for tables
            if slide.shapes:
                for shape in slide.shapes:
                    if shape.shape_type == 19:  # Table
                        f.write("\n[TABLE DETECTED]\n")
                        try:
                            for row in shape.table.rows:
                                row_text = " | ".join([cell.text for cell in row.cells])
                                f.write(row_text + "\n")
                        except:
                            pass
            
            f.write("-" * 80 + "\n")
        
        f.write("\n" + "="*80 + "\n")
        f.write("EXTRACTION COMPLETE\n")
        print("Output written to ppt_output.txt")
    
except Exception as e:
    print(f"Error: {e}")
    sys.exit(1)
